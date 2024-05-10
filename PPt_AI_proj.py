from pptx import Presentation
from openai import OpenAI
import os
import creds

def extract_powerpoint(pptx_path):

    prs = Presentation(pptx_path)
    pptx_text = ""
    for slide in prs.slides: 
        for shape in slide.shapes: 
            if hasattr(shape, "text"): 
                pptx_text += shape.text + "\n"
    return pptx_text

def gpt_for_lolz(pptx_text):
    client = OpenAI(
        api_key = creds.api_key
    )
    response = client.chat.completions.create(
        model = "gpt-3.5-turbo",
        messages = [
            {"role" : "system", "content" : "you are a helpful assistant that summarizes text"},
            {"role" : "user", "content" : 'please summarize this text with specific bullet points' + pptx_text}
        ]
    )
    return response.choices[0].message.content


def main():
    unstripped_pptx_path = input("paste path here")
    stripped_pptx_path = unstripped_pptx_path.strip('\'"')

    pptx_text = extract_powerpoint(stripped_pptx_path)
    output = gpt_for_lolz(pptx_text)
    f = open("summary.txt", "w")
    print(output, file = f)
    f.close()
    


main()

